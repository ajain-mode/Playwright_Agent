"""
Generate PowerPoint from agentic-system-slides.html content.
Output: docs/presentations/agentic-system-slides.pptx
"""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "agentic-system-slides.pptx"
DIAGRAM_PNG = ROOT / "agentic-pipeline-diagram.png"

PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
PRIMARY_LIGHT = RGBColor(0x2B, 0x6C, 0xB0)
ACCENT = RGBColor(0xED, 0x89, 0x36)
TEXT = RGBColor(0x2D, 0x37, 0x48)
TEXT_LIGHT = RGBColor(0x71, 0x80, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x38, 0xA1, 0x69)
SUCCESS_BG = RGBColor(0xF0, 0xFF, 0xF4)
IMPACT_BG = RGBColor(0xFF, 0xF5, 0xF5)
IMPACT_TEXT = RGBColor(0x74, 0x2A, 0x2A)


def capture_diagram() -> None:
    """Screenshot pipeline SVG from HTML via Playwright."""
    html_path = ROOT / "agentic-system-slides.html"
    script = f"""
const {{ chromium }} = require('playwright');
(async () => {{
  const browser = await chromium.launch({{ channel: 'chrome' }});
  const page = await browser.newPage({{ viewport: {{ width: 1000, height: 800 }} }});
  await page.goto('file:///{html_path.as_posix()}');
  const el = page.locator('.diagram-container');
  await el.waitFor({{ state: 'visible' }});
  await el.screenshot({{ path: '{DIAGRAM_PNG.as_posix()}', type: 'png' }});
  await browser.close();
}})();
"""
    tmp = ROOT / "_capture-diagram.js"
    tmp.write_text(script, encoding="utf-8")
    try:
        subprocess.run(
            ["node", str(tmp)],
            cwd=ROOT.parent.parent,
            check=True,
            capture_output=True,
            text=True,
        )
    except (subprocess.CalledProcessError, FileNotFoundError) as exc:
        print(f"Warning: diagram capture skipped ({exc})", file=sys.stderr)
    finally:
        if tmp.exists():
            tmp.unlink()


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    word_wrap: bool = True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items: list[str], font_size: int = 13):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)
    return box


def add_section_header(slide, num: str, title: str) -> None:
    circle = slide.shapes.add_shape(1, Inches(0.5), Inches(0.35), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(
        slide, Inches(1.05), Inches(0.32), Inches(8.5), Inches(0.55),
        title, font_size=26, bold=True, color=PRIMARY,
    )


def add_banner(slide, top, text: str, *, kind: str = "success") -> None:
    shape = slide.shapes.add_shape(1, Inches(0.5), top, Inches(9.0), Inches(0.85))
    shape.fill.solid()
    shape.line.fill.background()
    if kind == "success":
        shape.fill.fore_color.rgb = SUCCESS_BG
        color = SUCCESS
    else:
        shape.fill.fore_color.rgb = IMPACT_BG
        color = IMPACT_TEXT
    add_textbox(slide, Inches(0.65), top + Inches(0.1), Inches(8.7), Inches(0.7), text, font_size=12, color=color)


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]], font_size: int = 9):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), title, font_size=22, bold=True, color=PRIMARY)

    cols = len(headers)
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.35), Inches(0.9), Inches(9.3), Inches(5.8))
    table = table_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)


def build_presentation() -> Presentation:
    if not DIAGRAM_PNG.exists():
        capture_diagram()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ── Cover ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_textbox(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.35),
                "MODE GLOBAL · SUNTECK TMS · QA AI AGENT", font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(1.2),
                "Playwright Agentic Test Automation —\nCurrent State & Gaps",
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(8.6), Inches(0.9),
                "Hybrid rule-based + LLM multi-agent pipeline for SunTeck TMS — Claude-powered generation tiers, "
                "self-healing validation, structured test inputs, and Cursor-governed guardrails.",
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    stats = [("4", "Specialized Agents"), ("4", "AI Generation Calls"), ("28", "Scripts Generated"),
             ("130+", "Automation Patterns"), ("3", "Apps Covered")]
    x = 0.55
    for num, label in stats:
        box = slide.shapes.add_shape(1, Inches(x), Inches(3.2), Inches(1.65), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY_LIGHT
        box.line.color.rgb = ACCENT
        add_textbox(slide, Inches(x), Inches(3.25), Inches(1.65), Inches(0.45), num, font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(3.65), Inches(1.65), Inches(0.4), label, font_size=8, color=WHITE, align=PP_ALIGN.CENTER)
        x += 1.78

    meta = "Prepared By: R Systems International  |  Prepared For: Mode Global  |  Platform: SunTeck TMS (BTMS · DME · TNX)  |  Date: June 29, 2026"
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(9), Inches(0.35), meta, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # ── TOC ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.5), "Contents", font_size=28, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(2.5), [
        "1  What We Built — The Agent Pipeline",
        "2  Upstream Gaps & Next Steps",
    ], font_size=20)

    # ── Section 1 intro ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, "1", "What We Built — The Agent Pipeline")
    add_banner(slide, Inches(1.0),
               "Outcome: Written test cases flow through a four-agent pipeline and emerge as validated, runnable Playwright scripts.\n"
               "Delivered: 24 billing toggle and 4 DFB test scripts successfully generated through this architecture.")
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(0.8),
                "The pipeline combines deterministic automation rules (proven patterns, real UI locators) with Claude AI calls "
                "at four strategic points. Every script is reviewed by the Quality Guardian and Coverage Ensurer before it is saved.",
                font_size=13)

    # ── Agents ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.45), "Specialized Agents", font_size=24, bold=True, color=PRIMARY)
    agents = [
        ("✦ AI Script Writer", "Claude · 4 generation calls",
         "Writes and adapts test scripts. Claude CLI primary; Anthropic API backup."),
        ("① Step Interpreter", "Agent 1 · StepProcessor",
         "Understands each step — screen, tab, edit vs view, app context."),
        ("② Code Mapper", "Agent 2 · POMMethodMatcher",
         "Maps steps to automation using 130+ patterns and 5,900+ UI element IDs."),
        ("③ Quality Guardian", "Agent 3 · SpecValidator",
         "Validates scripts, auto-fixes violations, up to 2 correction passes."),
        ("④ Coverage Ensurer", "Agent 4 · SpecFeedbackLoop · 1 pass · up to 8 steps",
         "Compares script vs test case; fills missing steps; re-runs Quality Guardian if changed."),
    ]
    positions = [(0.4, 1.0), (5.1, 1.0), (0.4, 2.85), (5.1, 2.85), (2.75, 4.55)]
    for (alias, role, desc), (x, y) in zip(agents, positions):
        w = 4.5 if len(positions) == 5 and y > 4 else 4.5
        if alias.startswith("④"):
            w, x = 4.5, 2.75
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.35 if not alias.startswith("④") else 0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
        box.line.color.rgb = PRIMARY_LIGHT
        add_textbox(slide, Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.3), alias, font_size=12, bold=True, color=PRIMARY)
        add_textbox(slide, Inches(x + 0.1), Inches(y + 0.35), Inches(w - 0.2), Inches(0.22), role, font_size=9, color=TEXT_LIGHT)
        add_textbox(slide, Inches(x + 0.1), Inches(y + 0.55), Inches(w - 0.2), Inches(0.7), desc, font_size=10)

    # ── KPIs ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.45), "Key Metrics", font_size=24, bold=True, color=PRIMARY)
    kpis = [("24", "Billing toggle scripts generated"), ("4", "DFB scripts generated"),
            ("5,900+", "UI elements indexed"), ("25+", "Auto-fix quality rules")]
    x = 0.55
    for val, lbl in kpis:
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(2.05), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_LIGHT
        add_textbox(slide, Inches(x), Inches(1.7), Inches(2.05), Inches(0.5), val, font_size=28, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.1), Inches(2.3), Inches(1.85), Inches(0.6), lbl, font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        x += 2.25

    # ── Pipeline diagram ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.4),
                "End-to-End Pipeline — How a Test Case Becomes a Script", font_size=20, bold=True, color=PRIMARY)
    if DIAGRAM_PNG.exists():
        slide.shapes.add_picture(str(DIAGRAM_PNG), Inches(0.35), Inches(0.65), width=Inches(9.3))
    else:
        add_textbox(slide, Inches(0.5), Inches(2), Inches(9), Inches(1),
                    "[Pipeline diagram — open agentic-system-slides.html for full visual]",
                    font_size=14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # ── AI calls table ──
    add_table_slide(prs, "✦ AI Script Writer — Four Claude Calls",
                    ["Call", "Name", "What It Does", "When", "If Fails"],
                    [
                        ["1", "Data Enricher", "Fills missing test data fields", "Before generation when data incomplete", "Continues with partial data"],
                        ["2", "Full Script Writer", "Writes entire script using reference template", "First attempt (≤40 steps, reference exists)", "Falls back to Call 3 or rules"],
                        ["3", "Smart Adapter", "Clones reference; rewrites differing steps", "When ≥70% similar to existing test", "Keeps reference or Call 4"],
                        ["4", "Step Writer", "Writes code for unmatched individual steps", "Last resort after rules fail", "Placeholder for manual completion"],
                    ])

    # ── Agent roster ──
    add_table_slide(prs, "Agent Roster — Business Names & Responsibilities",
                    ["Business Name", "Component", "Role", "Output / Loops"],
                    [
                        ["AI Script Writer", "LLMService", "4 Claude calls for data & script generation", "Code → Script Assembler"],
                        ["Step Interpreter", "Agent 1", "Understands each written step", "Context → Code Mapper + Guardian"],
                        ["Code Mapper", "Agent 2", "Maps steps to proven actions", "Draft script; receives fixes from 3 & 4"],
                        ["Quality Guardian", "Agent 3", "Validates & auto-fixes scripts", "≤2 passes; loops to Code Mapper"],
                        ["Coverage Ensurer", "Agent 4", "Ensures every step in final script", "1 pass; max 8 steps; re-runs Guardian"],
                    ], font_size=8)

    # ── In / Out ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.45), "Inputs & Outputs", font_size=24, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(4.2), Inches(0.35), "What Goes In", font_size=16, bold=True, color=SUCCESS)
    add_bullets(slide, Inches(0.5), Inches(1.2), Inches(4.3), Inches(3.5), [
        "Written test cases — CSV, Excel, JSON, or plain text",
        "Step-by-step actions and expected results",
        "Runtime test data per category",
        "Application source (BTMS + DME repos) for UI locators",
    ])
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.2), Inches(0.35), "What Comes Out", font_size=16, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.5), [
        "Playwright .spec.ts scripts — 24 billing toggle & 4 DFB generated",
        "Organized by product category",
        "Traceable to original test case step numbers",
        "Parsed data saved to dedicated module data files",
    ])

    # ── Governance ──
    add_table_slide(prs, "Governance & Evidence — How Quality Is Enforced",
                    ["Artifact", "What It Ensures", "Type"],
                    [
                        ["Cursor governance rules", "Architecture standards, guardrails, POM protection", "Always on"],
                        ["130+ automation patterns", "Proven step-to-code mappings", "Rule-based"],
                        ["5,900+ indexed UI elements", "Real locators from app source", "App source"],
                        ["Reference test catalog", "Clone when ≥70% similar", "Pattern reuse"],
                        ["25+ auto-fix rules", "Quality Guardian corrects before review", "Self-healing"],
                    ], font_size=10)

    # ── Section 2 intro ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, "2", "Upstream Gaps & Next Steps")
    add_banner(slide, Inches(1.0),
               "The pipeline is built and working — but its effectiveness depends on the quality of inputs.",
               kind="impact")
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(0.6),
                "Thin test cases, undocumented business rules, or restricted application access will limit what the agents and AI Script Writer can produce.",
                font_size=13, color=IMPACT_TEXT)

    # ── Gaps 1 & 2 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(4.3), Inches(0.35), "① Test Cases Lack Sufficient Detail", font_size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(0.5), Inches(0.7), Inches(4.3), Inches(2.5), [
        "Testmo cases need full steps and expected outcomes",
        "Manual QA spends 3–5 hours per test case — largest upstream effort",
        "Without detail, spec generation is inaccurate",
    ], font_size=11)
    add_textbox(slide, Inches(5.1), Inches(0.3), Inches(4.3), Inches(0.35), "② Tribal Knowledge — Scenarios Without Steps", font_size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(5.1), Inches(0.7), Inches(4.3), Inches(2.5), [
        "Scenarios exist in JIRA — the what is documented",
        "Agent-ready step detail is often missing — the how is not",
        "Precondition paths cannot be automated reliably",
    ], font_size=11)

    # ── Gaps 3 & 4 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(4.3), Inches(0.35), "③ Limited Application Source Access", font_size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(0.5), Inches(0.7), Inches(4.3), Inches(2.5), [
        "UI Element Index syncs BTMS and DME repos only",
        "Code Mapper cannot discover new UI IDs without source",
        "Broader repo sync needed for full coverage",
    ], font_size=11)
    add_textbox(slide, Inches(5.1), Inches(0.3), Inches(4.3), Inches(0.35), "④ Execution Time — Monolithic Scripts", font_size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(5.1), Inches(0.7), Inches(4.3), Inches(2.5), [
        "Large cases become single long-running scripts",
        "No shared preconditions or phased build-up",
        "Multi-app flows multiply wall-clock time per run",
    ], font_size=11)

    # ── Next steps ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.45), "Recommended Next Steps", font_size=24, bold=True, color=SUCCESS)
    add_bullets(slide, Inches(0.6), Inches(1.0), Inches(4.2), Inches(4), [
        "Enrich test cases with explicit expected results per step",
        "Add agent-ready step detail to JIRA scenarios",
        "Approve additional application repo sync",
    ], font_size=13)
    add_bullets(slide, Inches(5.2), Inches(1.0), Inches(4.2), Inches(4), [
        "Decompose large cases into phased specs or shared fixtures",
        "Grow the reference test catalog per category",
        "Cap step count per script with continuation links",
    ], font_size=13)

    # ── Closing ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(0.6),
                "R Systems International  |  Confidential — Prepared for Mode Global",
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(0.4),
                "Agentic Test Automation Pipeline · SunTeck TMS · June 2026",
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUT_PATH)
    print(f"Created: {OUT_PATH}")


if __name__ == "__main__":
    main()
