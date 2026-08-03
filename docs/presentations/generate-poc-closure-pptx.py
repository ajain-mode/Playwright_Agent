"""
Generate 3-slide QA AI POC handover presentation for Mode Global.
Output: docs/presentations/POC-QA-Agent-Closure.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "POC-QA-Agent-Closure.pptx"
DIAGRAM = ROOT / "agentic-pipeline-diagram.png"
if not DIAGRAM.exists():
    DIAGRAM = ROOT / "pptx-v2-captures" / "07-diagram.png"

PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
PRIMARY_LIGHT = RGBColor(0x2B, 0x6C, 0xB0)
ACCENT = RGBColor(0xED, 0x89, 0x36)
TEXT = RGBColor(0x2D, 0x37, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x38, 0xA1, 0x69)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, left, top, width, height, text, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def bullets(slide, left, top, width, height, items: list[str], size=9, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(2)


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)
    textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), Inches(0.35),
            "MODE GLOBAL · SUNTECK TMS", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(0.9),
            "QA AI POC — Mode", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.5), Inches(2.35), Inches(9), Inches(0.55),
            "Playwright Agentic Test Automation — Handover & Deliverables", size=17, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.7), Inches(3.05), Inches(8.6), Inches(0.45),
            "Hybrid rule-based + LLM multi-agent pipeline · Claude-powered generation tiers · "
            "Self-healing validation · Cursor-governed guardrails",
            size=11, color=WHITE, align=PP_ALIGN.CENTER)

    stats = [("28", "Scripts Delivered"), ("24", "Billing Toggle"), ("4", "DFB"), ("5", "Agents + AI Writer")]
    x = 1.05
    for num, label in stats:
        box = slide.shapes.add_shape(1, Inches(x), Inches(3.75), Inches(1.7), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY_LIGHT
        box.line.color.rgb = ACCENT
        textbox(slide, Inches(x), Inches(3.8), Inches(1.7), Inches(0.4), num, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        textbox(slide, Inches(x), Inches(4.18), Inches(1.7), Inches(0.35), label, size=8, color=WHITE, align=PP_ALIGN.CENTER)
        x += 1.95

    textbox(slide, Inches(0.5), Inches(4.95), Inches(9), Inches(0.3),
            "R Systems International  |  Prepared for Mode Global  |  July 2026",
            size=10, color=WHITE, align=PP_ALIGN.CENTER)


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    textbox(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
            "Agentic Architecture — How a Test Case Becomes a Script", size=22, bold=True, color=PRIMARY)

    if DIAGRAM.exists():
        slide.shapes.add_picture(str(DIAGRAM), Inches(0.35), Inches(0.62), width=Inches(5.6))

    textbox(slide, Inches(6.05), Inches(0.62), Inches(3.55), Inches(0.28),
            "Pipeline Flow", size=11, bold=True, color=PRIMARY_LIGHT)
    bullets(slide, Inches(6.05), Inches(0.88), Inches(3.55), Inches(1.35), [
        "Input: CSV · Excel · JSON · plain text test cases",
        "Prep: Test Case Reader · App Context Sync · UI Element Index",
        "Script Assembler → Draft .spec.ts",
        "③ Quality Guardian (≤2 passes) → ④ Coverage Ensurer (1 pass)",
        "Output: Runnable Playwright test script",
    ], size=8)

    textbox(slide, Inches(6.05), Inches(2.2), Inches(3.55), Inches(0.25),
            "What Each Agent Does", size=11, bold=True, color=PRIMARY_LIGHT)
    bullets(slide, Inches(6.05), Inches(2.45), Inches(3.55), Inches(2.1), [
        "✦ AI Script Writer — 4 Claude calls: Data Enricher · Full Script Writer · Smart Adapter · Step Writer",
        "① Step Interpreter — screen, tab, edit/view, app context (BTMS/DME/TNX)",
        "② Code Mapper — 130+ patterns · 5,900+ UI IDs from app source",
        "③ Quality Guardian — 25+ auto-fix rules · guardrails · step coverage",
        "④ Coverage Ensurer — fills up to 8 missing steps · re-runs Guardian if changed",
    ], size=7.5)

    bar = slide.shapes.add_shape(1, Inches(0.35), Inches(4.72), Inches(9.3), Inches(0.78))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF4)
    bar.line.color.rgb = SUCCESS
    textbox(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.65),
            "Generation order: Call 2 → Call 3 → Rule-based (Step Interpreter + Code Mapper) → Call 4.  "
            "Claude CLI + Anthropic API · All AI output passes Quality Guardian & Coverage Ensurer before save.  "
            "POC delivered: 24 billing toggle + 4 DFB scripts.",
            size=9, color=RGBColor(0x22, 0x54, 0x3D))


def slide_deliverables(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "POC Deliverables — For Mode Global & QA Teams", size=24, bold=True, color=PRIMARY)

    textbox(slide, Inches(0.5), Inches(0.72), Inches(4.5), Inches(0.25),
            "What You Receive", size=13, bold=True, color=PRIMARY_LIGHT)
    bullets(slide, Inches(0.5), Inches(0.98), Inches(4.5), Inches(2.4), [
        "Playwright Agent repository — framework, page objects, agent pipeline",
        "src/tests/AIAgent/ — 28 generated .spec.ts scripts (billing toggle + DFB)",
        "src/data/<category>/ — runtime test data CSVs per product module",
        "src/agent/examples/sample-testcases.csv — testcase step catalog",
        "Agent CLI — generate scripts from detailed test cases",
        "Reference specs + Cursor governance rules",
        "Handover document + technical deck (agentic-system-slides.html)",
    ], size=9.5)

    textbox(slide, Inches(5.1), Inches(0.72), Inches(4.4), Inches(0.25),
            "How QA Teams Use the Agent Going Forward", size=13, bold=True, color=SUCCESS)
    bullets(slide, Inches(5.1), Inches(0.98), Inches(4.4), Inches(2.4), [
        "Run tests:  npx playwright test src/tests/AIAgent/",
        "Interactive generation:  npm run agent:cli",
        "From file (CSV/Excel/JSON):  npm run agent:file -- <path>",
        "Batch:  npm run agent:batch -- <folder>",
        "Preview only:  npm run agent:preview",
        "Prepare Testmo cases with full steps + expected results (3–5 hrs/case)",
        "Reports:  npm run test:allure  ·  npm run submit:testmo",
    ], size=9.5)

    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(4.45), Inches(9), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    textbox(slide, Inches(0.65), Inches(4.55), Inches(8.7), Inches(0.28),
            "Handover Complete", size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.65), Inches(4.85), Inches(8.7), Inches(0.5),
            "The agentic pipeline is built, validated, and ready for Mode Global QA. "
            "Run existing scripts today; generate new automation from detailed test cases via the Agent CLI. "
            "See QA-AI-Agent-Handover.docx for full setup, architecture, and step-by-step usage.",
            size=9.5, color=WHITE, align=PP_ALIGN.CENTER)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide_cover(prs)
    slide_architecture(prs)
    slide_deliverables(prs)
    prs.save(OUT_PATH)
    print(f"Created: {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
