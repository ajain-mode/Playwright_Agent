"""
Generate PowerPoint v2 from agentic-system-slides.html — exact HTML content via section screenshots.
Output: docs/presentations/agentic-system-slides-v2.pptx
Does NOT overwrite agentic-system-slides.pptx
"""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
HTML_PATH = ROOT / "agentic-system-slides.html"
CAPTURE_SCRIPT = ROOT / "capture-html-sections-v2.js"
CAPTURE_DIR = ROOT / "pptx-v2-captures"
OUT_PATH = ROOT / "agentic-system-slides-v2.pptx"

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def run_captures() -> None:
    repo_root = ROOT.parent.parent
    result = subprocess.run(
        ["node", str(CAPTURE_SCRIPT)],
        cwd=str(repo_root),
        capture_output=True,
        text=True,
    )
    if result.stdout:
        print(result.stdout.strip())
    if result.returncode != 0:
        print(result.stderr, file=sys.stderr)
        raise SystemExit(result.returncode)


def add_image_slide(prs: Presentation, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    img = Image.open(image_path)
    iw, ih = img.size
    aspect = iw / ih
    slide_aspect = float(slide_w) / float(slide_h)

    if aspect >= slide_aspect:
        pic_w = slide_w
        pic_h = int(slide_w / aspect)
    else:
        pic_h = slide_h
        pic_w = int(slide_h * aspect)

    left = int((slide_w - pic_w) / 2)
    top = int((slide_h - pic_h) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=pic_w, height=pic_h)


def build_presentation() -> Presentation:
    print("Running HTML section captures...")
    run_captures()

    images = sorted(
        p for p in CAPTURE_DIR.glob("*.png") if "_part" not in p.stem
    )
    if not images:
        raise FileNotFoundError(f"No captures found in {CAPTURE_DIR}.")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for img_path in images:
        add_image_slide(prs, img_path)

    return prs


def main() -> None:
    if not HTML_PATH.exists():
        raise FileNotFoundError(f"Missing HTML: {HTML_PATH}")
    prs = build_presentation()
    out = OUT_PATH
    try:
        prs.save(out)
    except PermissionError:
        out = ROOT / "agentic-system-slides-v2-generated.pptx"
        prs.save(out)
        print("Note: original v2 file was locked; saved alternate copy.", file=sys.stderr)
    print(f"Created: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
